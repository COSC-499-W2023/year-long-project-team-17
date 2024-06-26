from django.shortcuts import render, redirect
from django.contrib.auth import authenticate, login, logout
from django.views.generic import CreateView
from django.contrib import messages
from django.contrib.auth.models import Group, User
from django.urls import reverse_lazy
from django.views import generic
from .decorators import *
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
from .forms import SignUpForm, EditProfileForm, ProfilePageForm, ChangePasswordForm, EditProfilePageForm
from .models import Profile, Message, Presentations
from util.generate_summary import generate_summary
from util.generate_presentation import generate_presentation, get_presentation_info
from util.detect_plagiarism import detect_plagiarism
from util.generate_exercises import generate_exercises_from_prompt, generate_similar_exercises
from util.adapt_content import generate_adapted_content
from util.modify_presentation import check_user_message, generate_modification_assistant_response, generate_modified_presentation
from django.http import FileResponse
from django.core.cache import cache
from django.http import JsonResponse
from django.utils import timezone
from django.core.files.storage import FileSystemStorage
from threading import Thread
from uuid import uuid4

import subprocess
import platform
import docx
from pptx import Presentation
import PyPDF2
import logging
from django.http import HttpResponse, JsonResponse, Http404
from openai import OpenAI
import json
from django.contrib.auth.decorators import login_required
from django.db import models
from django_ratelimit.decorators import ratelimit
from django_ratelimit.core import get_usage, is_ratelimited
from humanfriendly import format_timespan
from django.db import connection
from django.http import JsonResponse
from .models import Message  


from django.core.mail import send_mail, BadHeaderError
from django.core.exceptions import ValidationError
from django.core.validators import validate_email
from util import config
from datetime import datetime, timezone
import math
from django.shortcuts import get_object_or_404
from django.core.paginator import Paginator
from collections import deque 
from django_htmx.http import HttpResponseClientRefresh

from reportlab.pdfgen import canvas

from django.http import HttpResponse
from django.conf import settings
import os


import subprocess
import platform
import os
from django.contrib.auth import update_session_auth_hash
# Create your views here.
# logging.info("----"*100)
#
# logging.info(os.environ.get("OPENAI_API_KEY"))
# print("----"*100)
#
# print(os.environ.get("OPENAI_API_KEY"))


client = OpenAI(
    api_key=os.environ.get("OPENAI_API_KEY"),
)

system_platform = platform.system()
if system_platform == 'Windows':
    presentation_root = 'media'
elif system_platform == 'Darwin':  # macOS
     presentation_root = 'app/media'
elif system_platform == 'Linux':
     presentation_root = 'app/media'
    
class AboutUsView(CreateView):
    template_name = 'about_page.html'
    model = Profile
    form_class = ProfilePageForm
    # fields = '__all__'


    def form_valid(self, form):
        form.instance.user = self.request.user
        return super().form_valid(form)
    

def upload(request):
    context = {}
    if request.method == 'POST':
        uploaded_file = request.FILES['document']
        fs = FileSystemStorage()
        name = fs.save(uploaded_file.name, uploaded_file)
        context['url'] = fs.url(name)
    return render(request, 'upload.html', context)



@login_required
def new_chats(request):
    users = User.objects.all().exclude(username=request.user.username)
    return render(request, 'new_chats.html', {'users': users})



def get_recent_messages(request, user_id):
    user_messages = Message.objects.filter(receiver=request.user).order_by('-timestamp')[:1]
    if user_messages.exists():
        latest_timestamp = user_messages[0].timestamp.date()  # Extracting only the date
        return JsonResponse({'timestamp': latest_timestamp})
    else:
        return JsonResponse({'timestamp': None})


@login_required
def chat(request, username):
    receiver = User.objects.get(username=username)
    messages = Message.objects.filter(
        (models.Q(sender=request.user) & models.Q(receiver=receiver)) |
        (models.Q(sender=receiver) & models.Q(receiver=request.user))
    ).order_by('timestamp')
    if request.htmx:
        return render(request, 'partials/chat_messages_load.html', {'messages' : messages, 'reciever': receiver})
    
    return render(request, 'chat.html', {'messages': messages, 'receiver': receiver})


@login_required
def send_message(request, username):
    if request.method == 'POST':
        receiver = User.objects.get(username=username)
        content = request.POST.get('content', '')

        if content:
            message = Message.objects.create(
                sender=request.user,
                receiver=receiver,
                content=content
            )

    return redirect('chat', username=username)


@login_required
def open_chats(request):
    # This querie will return the id of the chat partner, the last message time, and the id of the last message
    raw_query = """
        SELECT
    CASE
        WHEN sender_id = %s THEN receiver_id
        ELSE sender_id
    END AS chat_partner_id,
    MAX(timestamp) AS last_message_time,
    MAX(id) AS message_id
FROM
    website_message
WHERE
    sender_id = %s OR receiver_id = %s
GROUP BY
    chat_partner_id
ORDER BY
    last_message_time ASC;

    """

    # Execute the raw SQL query
    with connection.cursor() as cursor:
        cursor.execute(raw_query, [request.user.id, request.user.id, request.user.id])
        rows = cursor.fetchall()

    # Fetch Message objects based on the IDs returned by the query
    chat_ids = [row[2] for row in rows]
    # Get the message objects that have the same id as the ones returned by the query
    chats = Message.objects.filter(id__in=chat_ids).order_by('-id') 
    

    # print(rows)
    # print(chats)

    # Render the 'open_chats.html' template with the retrieved message objects
    
    return render(request, 'open_chats.html', {'chats': chats})



    
class CreateProfilePageView(CreateView):
    model = Profile
    form_class = ProfilePageForm
    template_name = 'create_user_profile_page.html'
    # fields = '__all__'

    def form_valid(self, form):
        form.instance.user = self.request.user
        return super().form_valid(form)


class EditProfilePageView(generic.UpdateView):
    model = Profile
    template_name = 'edit_profile_page.html'
    fields = ['bio', 'profile_pic', 'website_url', 'facebook_url', 'twitter_url', 'instagram_url', 'pinterest_url']
    success_url = reverse_lazy('home')


class UserEditView(generic.UpdateView):
    form_class = EditProfileForm
    template_name = 'edit_profile.html'
    success_url = reverse_lazy('home')

    def get_object(self):
        return self.request.user


def isLimited(request, exception):
    values = {}
    values['timeleft'] = 0

    if request.get_full_path() == "/":
        usage = get_usage(request, method=get_usage.ALL, key='post:username',
                       rate='5/5m', group='a')
        values['timeleft'] = format_timespan(usage['time_left'])
        return render(request, 'limited.html', values, status=429)
      
    elif request.get_full_path() == "/reset_password/":
        usage = get_usage(request, method=get_usage.ALL, key='ip', rate='8/5h', group='b')
        values['timeleft'] = format_timespan(usage['time_left'])
        return render(request, 'limited.html', values, status=429)


def faq(request):
    return render(request,'faq.html')


def contact_us(request):
    if request.method == 'POST':
        email = request.POST['email']
        message = request.POST['question']
        if(len(message) == 0):
            messages.error(request, "Please enter a question.")
            return render(request, 'contact_us.html')
        try:
            validate_email(email)
            subject = 'Question from ' + email   
            #Sends question from user as email to noreplyeduprompt@gmail.com
            send_mail(subject,
                    message,
                    settings.EMAIL_HOST_USER,
                    ['noreplyeduprompt@gmail.com'],
                    fail_silently=False)
            messages.success(request, "Your question has been sent.")
            return render(request, 'contact_us.html')
        except ValidationError as err:
            messages.error(request, "An invalid email address was entered please try again.")
            return render(request, 'contact_us.html')
        except BadHeaderError as err:
            messages.error(request, "An invalid header was found please try again.")
            return render(request, 'contact_us.html')
    else:
        return render(request, 'contact_us.html')


@authenticated_user
def Profile(request, username):

    #Users viewing their own profile page
    if request.user.username == username:
        results = Presentations.objects.filter(user_id = request.user.id).values('main_title','titles','date_created', 'is_shared').order_by('-date_created')
        context = {}
        formated_date = []
        for value in results.values():
            time = datetime.now(timezone.utc) - value['date_created']
            time = format_timespan(math.floor(time.total_seconds()), max_units=2)
            formated_date.append(time)

        page_results_number = int(request.GET.get('page_results', 1))
        paginator_results = Paginator(results, 2)
        page_results_obj = paginator_results.get_page(page_results_number)
        
        page_date_number = int(request.GET.get('page_date', 1))
        paginator_date = Paginator(formated_date, 2)
        page_date_obj = paginator_date.get_page(page_date_number)
        
        context['page_results_obj'] = page_results_obj
        page_date_obj.object_list = deque(page_date_obj.object_list)
        context['page_date_obj'] = page_date_obj
        
        if request.htmx:
            return render(request, 'partials/profile_list.html', context)

        return render(request, 'profile.html', context)
    else:
        #Users viewing a users profile
        context = {}
        #if username does not exist raise 404 not found http exception
        user_obj = get_object_or_404(User, username = username)
        #Get presentation that are public
        results = Presentations.objects.filter(user_id = user_obj.pk, is_shared = 1).values('main_title','titles','date_created').order_by('-date_created')
        formated_date = []
        for value in results.values():
            time = datetime.now(timezone.utc) - value['date_created']
            time = format_timespan(math.floor(time.total_seconds()), max_units = 2)
            formated_date.append(time)
        
        

        page_results_number = int(request.GET.get('page_results', 1))
        paginator_results = Paginator(results, 2)
        page_results_obj = paginator_results.get_page(page_results_number)

        page_date_number = int(request.GET.get('page_date', 1))
        paginator_date = Paginator(formated_date, 2)
        page_date_obj = paginator_date.get_page(page_date_number)
        
        context['user_obj'] = user_obj
        context['page_results_obj'] = page_results_obj
        page_date_obj.object_list = deque(page_date_obj.object_list)
        context['page_date_obj'] = page_date_obj
        if request.htmx:
            return render(request, 'partials/profile_different_user_list.html', context)

        return render(request, 'profile_different_user.html', context)


@authenticated_user
def forumPage(request):
    #Users viewing a users profile
    context = {}
    #if username does not exist raise 404 not found http exception
    #Get presentation that are public
    results = Presentations.objects.filter(is_shared = 1).values('main_title','titles','date_created').order_by('-date_created')
    
    formated_results = []
    for value in results.values():
        time = datetime.now(timezone.utc) - value['date_created']
        time = format_timespan(math.floor(time.total_seconds()), max_units = 2)
        
        
        user = User.objects.get(id=value['user_id'])
        username = user.username

        formated_results.append({
            'id': value['id'],
            'user_id': value['user_id'],
            'username': username,
            'is_shared': value['is_shared'],
            'main_title': value['main_title'],
            'titles': value['titles'],
            'date_created': value['date_created'],
            'formatted_time': time
        })
    

    if request.htmx:
        return render(request, 'partials/profile_different_user_list_forum.html', context)

    return render(request, 'forumPage.html', {'values' : formated_results})

from django.db.models import Q

@authenticated_user
def search_results(request):
    context = {}
    query = request.GET.get('query')
    results = Presentations.objects.filter(
        Q(main_title__icontains=query) | Q(titles__icontains=query), is_shared = 1
    )
    formated_results = []
    for value in results.values():
        time = datetime.now(timezone.utc) - value['date_created']
        time = format_timespan(math.floor(time.total_seconds()), max_units = 2)

    
        user = User.objects.get(id=value['user_id'])
        username = user.username

        formated_results.append({
            'id': value['id'],
            'user_id': value['user_id'],
            'username': username,
            'is_shared': value['is_shared'],
            'main_title': value['main_title'],
            'titles': value['titles'],
            'date_created': value['date_created'],
            'formatted_time': time
        })
    if request.htmx:
        return render(request, 'partials/profile_different_user_list_forum.html', context)
    if( query == ""):
        return render(request, 'forumPage.html', {'values' : formated_results})
    return render(request, 'search_results.html', {'results': formated_results, 'query': query})


@authenticated_user
def download_presentation_pptx(request, pres_id):
    #raises 404 http exception if presentation object does not exist
    pres = get_object_or_404(Presentations, pk=pres_id)
    if(pres.is_shared == 1 or pres.user_id == request.user.id):
        filename = pres.presentation
        base_dir = str(settings.BASE_DIR)
        file_path = base_dir + "/media/presentations/" + filename
        try:
            pres = Presentation(file_path)
            response = HttpResponse(
            content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
            response['Content-Disposition'] = f'attachment; filename="{filename}"'
            pres.save(response)
            return response
        except:
            raise Http404("File not found.")
    else:
        messages.error(request, "Sorry this presentation has been set to private and can no longer be downloaded.")
        return redirect("home")


@authenticated_user
def change_post_visibility(request, pres_id):
    if request.method == 'POST' and request.htmx:
        if request.headers.get('HX-Trigger') == 'preview-pres' or request.headers.get('HX-Trigger') == 'post-visibility':
            #raises 404 http exception if presentation object does not exist
            pres = get_object_or_404(Presentations, pk=pres_id)
            #only update if presentation belongs to user making request
            if request.user.id == pres.user_id:  
                if pres.is_shared == 0:  
                    pres.is_shared = 1
                else:
                    pres.is_shared = 0
                pres.save(update_fields=["is_shared"])
                if request.headers.get('HX-Trigger') == 'preview-pres':
                    return render(request, 'partials/preview_post_visibility.html', {'presentation_obj': pres})
                elif request.headers.get('HX-Trigger') == 'post-visibility':
                    return render(request, 'partials/post_visibility.html', {'value': pres})
            else:
                messages.error(request, "You do not have the authorization to make changes to that post.")
                return redirect("home")
        else:  
            messages.error(request, "You do not have the authorization to make changes to that post.")
            return redirect("home")
    else:  
        messages.error(request, "You do not have the authorization to make changes to that post.")
        return redirect("home")


@authenticated_user
def delete_presentation(request, pres_id):
    if request.htmx:
        #raises 404 http exception if presentation object does not exist
        pres = get_object_or_404(Presentations, pk=pres_id)
        #only delete if presentation belongs to user making request
        if request.user.id == pres.user_id:
            fs = FileSystemStorage(location=presentation_root + "/presentations/")    
            #Remove presentation from storage
            if fs.exists(pres.presentation):
                fs.delete(pres.presentation)
            pres.delete()
            messages.success(request, "Your presentation has been deleted.")
            return HttpResponseClientRefresh()
        else:
            messages.error(request, "You do not have the authorization to make changes to that post.")
            return redirect("home")      
    else:
        messages.error(request, "You do not have the authorization to make changes to that post.")
        return redirect("home")
    

@ratelimit(key='post:username', rate='5/5m',
           method=['POST'], group='a')
def home(request):


    #check to see if logging in
    if request.method=="POST":
        username = request.POST["username"]
        password = request.POST["password"]
        #authenticate
        user = authenticate(request, username=username, password=password)

        if user is not None:
            login(request, user)
            messages.success(request, "You have been Logged In")
            return redirect('home')
        else:
            messages.error(request, "There was an error logging in, please try again...")
            return redirect('home')
        
    else:
        return render(request, 'home.html')


@authenticated_user
def logout_user(request):
    logout(request)
    messages.success(request, "You have been Logged Out...")
    return redirect('home')


@unauthenticated_user
def register_user(request):
    if request.method == 'POST':
        form = SignUpForm(request.POST)
        if form.is_valid():
            user = form.save()
            #Authenticate and login
            username = form.cleaned_data['username']
            password = form.cleaned_data['password1']
            groupName = form.cleaned_data['user_group']
            group = Group.objects.get(name=groupName)
            user = authenticate(username = username, password = password)
            user.groups.add(group)
            login(request, user)
            messages.success(request, "You have succesfully registered")
            return redirect('home')
        else: 
            return render(request, 'register.html', {'form':form})
    else:
        form = SignUpForm()
        return render(request, 'register.html', {'form':form})

@authenticated_user
def change_password(request):
        if request.method == 'POST':    
            form = ChangePasswordForm(user=request.user, data=request.POST)
            if form.is_valid():
                form.save()
                #Updates session with the new session hash after password change so that the user stays logged in
                update_session_auth_hash(request, request.user)
                messages.success(request, "Your password has been changed successfully.")
                form = ChangePasswordForm(user=request.user)
                return render(request, 'change_password.html', {'form': form})
            else:
                return render(request,'change_password.html',{'form': form})
        else:
            form = ChangePasswordForm(user=request.user)
            return render(request, 'change_password.html', {'form': form})
    
@authenticated_user
def edit_profile(request):
    if request.method == 'POST':
        if 'settings_form' in request.POST:
            settings_form = EditProfileForm(data=request.POST, instance=request.user)
            if settings_form.is_valid():
                settings_form.save()
                page_form = EditProfilePageForm(initial={'bio' : request.user.profile.bio})
                settings_form = EditProfileForm(instance=request.user)
                messages.success(request, "Your information has been updated.")
                return render(request, 'edit_profile.html', {'page_form' : page_form, 'settings_form' : settings_form})
            else:
                page_form = EditProfilePageForm(initial={'bio' : request.user.profile.bio})
                messages.error(request, "A problem occured when updating your settings, your form contains the following errors listed below.")
                return render(request, 'edit_profile.html', {'page_form' : page_form, 'settings_form' : settings_form})
        if 'page_form' in request.POST:
            page_form = EditProfilePageForm(data = request.POST, files=request.FILES, instance = request.user.profile)
            if page_form.is_valid():
                page_form.save()
                page_form = EditProfilePageForm(initial={'bio' : request.user.profile.bio})
                settings_form = EditProfileForm(instance=request.user)
                messages.success(request, "Your profile page information has been updated.")
                return render(request, 'edit_profile.html', {'page_form' : page_form, 'settings_form' : settings_form})
            else:
                settings_form = EditProfileForm(instance=request.user)
                messages.error(request, "A problem occurred when updating your profile page information, your form contains the following errors listed below.")
                return render(request, 'edit_profile.html', {'page_form' : page_form, 'settings_form' : settings_form})
        else:
            page_form = EditProfilePageForm(initial={'bio' : request.user.profile.bio})
            settings_form = EditProfileForm(instance=request.user)
            messages.error(request, "There was a problem updating your information, please try again.")
            return render(request, 'edit_profile.html', {'page_form' : page_form, 'settings_form' : settings_form})
    else:
        page_form = EditProfilePageForm(initial={'bio' : request.user.profile.bio})
        settings_form = EditProfileForm(instance=request.user)
        return render(request, 'edit_profile.html', {'page_form' : page_form, 'settings_form' : settings_form})


#DELETE
def customer_record(request, pk):
    if request.user.is_authenticated:
        #look up record
        customer_record = Record.objects.get(id=pk)
        return render(request, 'record.html', {'customer_record':customer_record})
    else:
        messages.success(request, "You must be logged in to view that page")
        return redirect('home')


#DELETE
@allowed_users(allowed_roles=['teacher'])
def delete_record(request, pk):
    if request.user.is_authenticated:
        delete_it = Record.objects.get(id=pk)
        delete_it.delete()
        messages.success(request, "The customer has been deleted succesfully")
        return redirect('home')
    else:
        messages.success(request, "You must be logged in do that action")
        return redirect('home')


#DELETE
@allowed_users(allowed_roles=['teacher'])
def add_record(request):
    form = AddRecordForm(request.POST or None)
    if request.user.is_authenticated:
        if request.method == "POST":
            if form.is_valid():
                add_record = form.save()
                messages.success(request, "Record Added...")
                return redirect('home')

        return render(request, 'add_record.html', {'form':form})
    else:
        messages.success(request, "You must be logged in")
        return redirect('home')


#DELETE
@allowed_users(allowed_roles=['teacher'])
def update_record(request, pk):
    if request.user.is_authenticated:
        current_record = Record.objects.get(id=pk)
        form = AddRecordForm(request.POST or None, instance = current_record)
        if form.is_valid():
            form.save()
            messages.success(request, "Record has been updated")
            return redirect('home')
        return render(request, 'update_record.html', {'form':form})
    else:
        messages.success(request, "You must be logged in")
        return redirect('home')


@authenticated_user
def generate_summary_view(request):
    """
    A view to generate a summary based on the user input
    :param request: The request containing user message
    :return: Generated summary
    """
    
    if request.method == "POST":
        input_option = request.POST.get("input_option")
        if input_option == "write":

            input_text = request.POST.get("input_text")
            if input_text:
                summary = generate_summary(input_text)
                return render(request, 'summary_generation.html',
                                {'summary': summary,
                                "input_option": input_option})
            else:
                messages.error(request, "Please enter a text that you want to summarize")
                return render(request, "summary_generation.html", {"input_option": "write"})
        elif input_option == "upload":
            uploaded_file = request.FILES.get("file")

            if uploaded_file:
                logging.info("Upload")
                file_extension = os.path.splitext(uploaded_file.name)[1].lower()
                print(file_extension)
                if file_extension in [".txt", ".pptx", ".docx", ".pdf"]:
                    if file_extension == ".docx":
                        doc = docx.Document(uploaded_file)
                        full_text = []
                        for paragraph in doc.paragraphs:
                            full_text.append(paragraph.text)

                        document_text = "\n".join(full_text)
                        if document_text:
                            summary = generate_summary(document_text)
                            return render(request, "summary_generation.html",
                                            {"summary": summary,
                                            "input_option": input_option})
                        else:
                            messages.error(request,
                                            "You uploaded an empty .docx file.", {"input_option": "upload"})

                    elif file_extension == ".pptx":
                        presentation = Presentation(uploaded_file)

                        slide_text = []

                        for slide in presentation.slides:
                            slide_text.append("")
                            for shape in slide.shapes:
                                if hasattr(shape, "text"):
                                    slide_text.append(shape.text)

                        presentation_text = "\n".join(slide_text)
                        if presentation_text:
                            summary = generate_summary(presentation_text)
                            return render(request, "summary_generation.html",
                                            {"summary": summary,
                                            "input_option": input_option})
                        else:
                            messages.error(request,
                                            "You uploaded an empty .pptx file.", {"input_option": "upload"})

                    elif file_extension == ".txt":
                        txt_text = uploaded_file.read().decode('utf-8')
                        if txt_text:
                            summary = generate_summary(txt_text)
                            return render(request, "summary_generation.html",
                                            {"summary": summary,
                                            "input_option": input_option})
                        else:
                            messages.error(request,
                                            "You uploaded an empty .txt file.", {"input_option": "upload"})

                    elif file_extension == ".pdf":
                        pdf_content = uploaded_file
                        pdf_file = PyPDF2.PdfReader(pdf_content)
                        pdf_text = ""

                        for page_num in range(len(pdf_file.pages)):
                            page = pdf_file.pages[page_num]
                            pdf_text += page.extract_text()

                        if pdf_text:
                            summary = generate_summary(pdf_text)
                            return render(request, "summary_generation.html",
                                            {"summary": summary,
                                            "input_option": input_option})
                        else:
                            messages.error(request,
                                            "You uploaded an empty .pdf file.", {"input_option": "upload"})
                else:
                    messages.error(request,
                                    "The Uploaded file must have one of the following extensions: .docs, .pptx, .txt, .pdf", {"input_option": "upload"})

            else:
                messages.error(request, "Please upload a file")
                render(request, "summary_generation.html", {"input_option": "upload"})
    return render(request, "summary_generation.html", {"input_option": "write"})


# def pptx_to_pdf(pptx_file, pdf_file):
#     # Load the PowerPoint presentation
#     prs = Presentation(pptx_file)
#
#     # Create a canvas with PDF size
#     c = canvas.Canvas(pdf_file, pagesize=letter)
#
#     # Set up dimensions
#     width, height = letter
#     slide_width = width - 72  # 1 inch margin on each side
#     slide_height = height - 72
#
#     # Go through each slide in the presentation
#     for slide in prs.slides:
#         # Start a new page for each slide
#         c.showPage()
#
#         # Draw the slide onto the canvas
#         slide.shapes._spTree.draw(c, 0, 0, slide_width, slide_height)
#
#     # Save the PDF
#     c.save()

def generate_new_file_id():
    """
    A function to generate a new file uuid which will be used to store generated presentation
    :return: Generated file id
    """
    current_uuid = uuid4()
    cache.set("generated_presentation_id", str(current_uuid), timeout=1500)
    return current_uuid

def convert_pptx_to_pdf(input_path, output_path=None):
    """
    A function that converts from pptx to pdf
    :param input_path: The file path to pptx
    :param output_path: The file path where pdf has to be stored
    :return: None (pdf is being stored of successfully executed)
    """

    system_platform = platform.system()
    if system_platform == 'Windows':
        libreoffice_path = 'C:/Program Files/LibreOffice/program/soffice.exe'
    elif system_platform == 'Darwin':  # macOS
        libreoffice_path = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
    elif system_platform == 'Linux':
        libreoffice_path = '/usr/bin/soffice'  # Adjust this path according to your Linux distribution
    else:
        print("Unsupported operating system.")
        return

    if output_path is None:
        output_path = os.getcwd()
        output_path = presentation_root + "/presentations/"
        print(output_path)

    command = [libreoffice_path, '--convert-to', 'pdf', '--outdir', output_path, input_path]

    try:
        subprocess.run(command, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        print(f"Conversion successful: '{input_path}' to PDF.")
    except Exception as e:
        print(f"Error converting file: {e}")


def presentation_generation_task(request, input_text):
    """
    A function to generate a presentation using user input
    :param request: Request containing user input message
    :param input_text: The user input
    :return: None (Presentation is generated)
    """
    if "generated_presentation_filename" in request.session:
        request.session['generated_presentation_filename'] = None

    fs = FileSystemStorage(location=presentation_root + "/presentations/")
    new_id = generate_new_file_id()
    filename = f'generated_presentation_{new_id}.pptx'  # Choose a unique filename if necessary
    if fs.exists(filename):
        fs.delete(filename)
    #presentation = generate_presentation(input_text)
    values, presentation_json = generate_presentation(input_text)
    presentation = values['presentation']
    pres_info = values['pres_info']
    pres = Presentations.objects.create(
        user=request.user,
        main_title=pres_info['main_title'],
        titles=pres_info['titles'],
        presentation=filename
    )
    cache.set("generated_presentation_json", json.dumps(presentation_json), timeout=1500)
    # with fs.open(filename, 'wb') as pptx_file:
    full_file_name = presentation_root + "/presentations/" + filename
    presentation.save(full_file_name)
    request.session['generated_presentation_filename'] = filename
    pdf_filename = f'generated_presentation_{new_id}.pdf'

    if fs.exists(filename):
        pptx_file_path = fs.path(filename)
        # pdf_file_path = fs.path(pdf_filename)
        # Popen(['unoconv', '-f', 'pdf', '-o', pdf_file_path, ppt_file_path]) #`sudo apt-get install -y unoconv
        # while not os.path.exists(pdf_file_path):
        #     time.sleep(1)
        # pptx_to_pdf(pptx_file_path, pdf_file_path)

        convert_pptx_to_pdf(pptx_file_path)


@authenticated_user
def generate_presentation_view(request):
    """
    A view to generate a presentation
    :param request: Request containing user message
    :return: Generated presentation
    """

    if request.method == "POST":
        input_option = request.POST.get("input_option")
        if input_option == "write":
            input_text = request.POST.get("input_text")
            if input_text:
                thread = Thread(target=presentation_generation_task, args=(request, input_text))
                thread.start()

                return JsonResponse({'status': 'success', 'message': 'Presentation generated'})
            else:
                messages.error(request, "Please enter text for the presentation")
                return render(request, "presentation_generation.html", {"input_option": "write"})
        elif input_option == "upload":
            uploaded_file = request.FILES.get("file")
            if uploaded_file:
                file_extension = os.path.splitext(uploaded_file.name)[1].lower()

                if file_extension in [".txt", ".docx", ".pdf"]:
                    if file_extension == ".docx":
                        doc = docx.Document(uploaded_file)
                        full_text = []
                        for paragraph in doc.paragraphs:
                            full_text.append(paragraph.text)

                        document_text = "\n".join(full_text)
                        if document_text:
                            # presentation = generate_presentation()
                            # response = HttpResponse(
                            #     content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
                            # response['Content-Disposition'] = 'attachment; filename="generated_presentation.pptx"'
                            # presentation.save(response)
                            # return response
                            input_text = "I want a presentation based on the following content:\n" + str(document_text)
                            thread = Thread(target=presentation_generation_task, args=(request, input_text))
                            thread.start()

                            return JsonResponse({'status': 'success', 'message': 'Presentation generated'})
                        else:
                            messages.error(request, "You uploaded an empty .docx file.", {"input_option": "upload"})

                    elif file_extension == ".txt":
                        txt_text = uploaded_file.read().decode('utf-8')
                        if txt_text:
                            # presentation = generate_presentation(
                            #     "I want a presentation based on the following content:\n" + str(txt_text))
                            # response = HttpResponse(
                            #     content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
                            # response['Content-Disposition'] = 'attachment; filename="generated_presentation.pptx"'
                            # presentation.save(response)
                            # return response
                            input_text = "I want a presentation based on the following content:\n" + str(txt_text)
                            thread = Thread(target=presentation_generation_task, args=(request, input_text))
                            thread.start()

                            return JsonResponse({'status': 'success', 'message': 'Presentation generated'})
                        else:
                            messages.error(request,
                                            "You uploaded an empty .txt file.", {"input_option": "upload"})

                    elif file_extension == ".pdf":
                        pdf_content = uploaded_file
                        pdf_file = PyPDF2.PdfReader(pdf_content)
                        pdf_text = ""

                        for page_num in range(len(pdf_file.pages)):
                            page = pdf_file.pages[page_num]
                            pdf_text += page.extract_text()

                        if pdf_text:
                            # presentation = generate_presentation(
                            #     "I want a presentation based on the following content:\n" + str(pdf_text))
                            # response = HttpResponse(
                            #     content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
                            # response['Content-Disposition'] = 'attachment; filename="generated_presentation.pptx"'
                            # presentation.save(response)
                            # return response
                            input_text = "I want a presentation based on the following content:\n" + str(pdf_text)
                            thread = Thread(target=presentation_generation_task, args=(request, input_text))
                            thread.start()

                            return JsonResponse({'status': 'success', 'message': 'Presentation generated'})
                        else:
                            messages.error(request,
                                            "You uploaded an empty .pdf file.", {"input_option": "upload"})
                else:
                    messages.error(request,
                                    "The Uploaded file must have one of the following extensions: .docs, .pptx, .txt, .pdf", {"input_option": "upload"})

            else:
                messages.error(request, "Please upload a file", {"input_option": "upload"})

    return render(request, "presentation_generation.html", {"input_option": "write"})
    

@authenticated_user
@allowed_users(allowed_roles=['teacher'])
def detect_plagiarism_view(request):
    """
    A view to detect plagiarism
    :param request: Request containing user message
    :return: Similarity scores between each pair of files
    """
    if request.method == "POST":
        uploaded_files = request.FILES.getlist('file')
        if uploaded_files and len(uploaded_files) > 1:
            for uploaded_file in uploaded_files:
                file_extension = os.path.splitext(uploaded_file.name)[1].lower()
                if file_extension not in [".docx", ".pdf"]:
                    messages.error(request, "The uploaded files must have one of the following extensions: .docx, .pdf")
                    return render(request, "plagiarism_detection.html")
            uploaded_files_names = []
            uploaded_files_texts = []
            for uploaded_file in uploaded_files:
                uploaded_files_names.append(uploaded_file.name)
                file_extension = os.path.splitext(uploaded_file.name)[1].lower()
                if file_extension == ".docx":
                    doc = docx.Document(uploaded_file)
                    full_text = []
                    for paragraph in doc.paragraphs:
                        full_text.append(paragraph.text)

                    document_text = "\n".join(full_text)
                    if document_text.strip():
                        uploaded_files_texts.append(document_text)
                        print(document_text)
                    else:
                        messages.error(request,
                                        "You uploaded an empty .docx file. The file name is the following: " + uploaded_file.name)
                        return render(request, "plagiarism_detection.html")

                elif file_extension == ".pdf":
                    pdf_content = uploaded_file
                    pdf_file = PyPDF2.PdfReader(pdf_content)
                    pdf_text = ""

                    for page_num in range(len(pdf_file.pages)):
                        page = pdf_file.pages[page_num]
                        pdf_text += page.extract_text()

                    if pdf_text.strip():
                        uploaded_files_texts.append(pdf_text)
                        print(pdf_text)
                    else:
                        messages.error(request,
                                        "You uploaded an empty .pdf file. The file name is the following: " + uploaded_file.name)
                        return render(request, "plagiarism_detection.html")

            plagiarised_files = []
            similarity_for_each_file_pair = []
            if uploaded_files_texts:
                plagiarism_pairs, all_similarity_scores = detect_plagiarism(uploaded_files_texts, 0.85)
                print(plagiarism_pairs)
                print(all_similarity_scores)
                if plagiarism_pairs:
                    for plag_pairs in plagiarism_pairs:
                        plagiarised_files.append((uploaded_files_names[plag_pairs[0]], uploaded_files_names[plag_pairs[1]], plag_pairs[2]*100))
                if all_similarity_scores:
                    for sim_score in all_similarity_scores:
                        similarity_for_each_file_pair.append((uploaded_files_names[sim_score[0]], uploaded_files_names[sim_score[1]], sim_score[2]*100))
            if plagiarised_files:
                print(plagiarised_files)
            if similarity_for_each_file_pair:
                print(similarity_for_each_file_pair)

            return render(request, "plagiarism_detection.html",
                            {"plagiarised_files": plagiarised_files, "all_similarities": similarity_for_each_file_pair})

        elif len(uploaded_files) == 1:
            messages.error(request, "Please upload more than one file.")
            return render(request, "plagiarism_detection.html")
        else:
            messages.error(request, "Please upload files before clicking the button. The number of files should be at least two.")
            return render(request, "plagiarism_detection.html")
    return render(request, "plagiarism_detection.html")


def generate_exercise_file_id():
    """
    A function to generate a new file uuid which will be used to store generated exercises
    :return: Generated file id
    """
    current_uuid = uuid4()
    cache.set("generated_exercise_id", str(current_uuid), timeout=1500)
    return current_uuid


from docx import Document


def exercise_generation_task(request, input_text, input_option):
    """
   A function to generate an exercise using user input
   :param request: Request containing user input message
   :param input_text: The user input
   :return: None (Exercise is generated)
   """
    if "generated_exercise_filename" in request.session:
        request.session['generated_exercise_filename'] = None

    fs = FileSystemStorage()
    new_id = generate_exercise_file_id()
    filename = f'generated_exercise_{new_id}.docx'  # Choose a unique filename if necessary
    if fs.exists(filename):
        fs.delete(filename)

    generated_exercises = generate_exercises_from_prompt(input_text)
    generated_exercise_doc = Document()

    generated_exercise_doc.add_paragraph(generated_exercises)
    with fs.open(filename, 'wb') as docx_file:
        generated_exercise_doc.save(docx_file)
    request.session['generated_exercise_filename'] = filename
    cache.set("input_option", input_option, timeout=1500)
    # request.session['input_option'] = input_option


@authenticated_user
def generate_exercise_view(request):
    """
   A view to generate exercises
   :param request: Request containing user message
   :return: Generated exercises
   """
    if request.method == "POST":
        input_option = request.POST.get("input_option")
        if input_option == "write":
            input_text = request.POST.get("input_text")
            if input_text:

                thread = Thread(target=exercise_generation_task, args=(request, input_text, input_option))
                thread.start()

                return JsonResponse({'status': 'success', 'message': 'Exercise generated'})
            else:
                messages.error(request, "Please describe what type of exercises do you want.")
                return render(request, "exercise_generation.html", {"input_option": "write"})
        elif input_option == "upload":
            uploaded_file = request.FILES.get("file")
            if uploaded_file:
                file_extension = os.path.splitext(uploaded_file.name)[1].lower()
                if file_extension not in [".docx", ".pdf"]:
                    messages.error(request,
                                    "The uploaded file must have one of the following extensions: .docx, .pdf")
                    return render(request, "exercise_generation.html", {"input_option": "upload"})
                elif file_extension == ".docx":
                        doc = docx.Document(uploaded_file)
                        full_text = []
                        for paragraph in doc.paragraphs:
                            full_text.append(paragraph.text)

                        document_text = "\n".join(full_text)
                        if document_text.strip():
                            thread = Thread(target=exercise_generation_task,
                                            args=(request, document_text, input_option))
                            thread.start()

                            return JsonResponse({'status': 'success', 'message': 'Exercise generated'})

                        else:
                            messages.error(request,
                                            "You uploaded an empty .docx file.")
                elif file_extension == ".pdf":
                    pdf_content = uploaded_file
                    pdf_file = PyPDF2.PdfReader(pdf_content)
                    pdf_text = ""

                    for page_num in range(len(pdf_file.pages)):
                        page = pdf_file.pages[page_num]
                        pdf_text += page.extract_text()

                    if pdf_text.strip():
                        thread = Thread(target=exercise_generation_task,
                                        args=(request, pdf_text, input_option))
                        thread.start()

                        return JsonResponse({'status': 'success', 'message': 'Exercise generated'})

                    else:
                        messages.error(request,
                                        "You uploaded an empty .pdf file.")
            else:
                messages.error(request,
                                "Please upload a file before Pressing the button. The uploaded file must have one of the following extensions: .docx, .pdf")
                return render(request, "exercise_generation.html", {"input_option": "upload"})
    return render(request, "exercise_generation.html", {"input_option": "write"})


@authenticated_user
def chatbot_view(request):
    """
   A view to generate the chatbot response to the user message
   :param request: Request containing user message
   :return: Generated chatbot response
   """
    chat_history = request.session.get('chat_history', [])

    if request.method == "POST":
        message = request.POST.get('message')
        chat_history_string = ""
        if chat_history:
            chat_history_string = "\n".join([f'User: {entry.get("USER")}\nAI: {entry.get("AI")}' for entry in chat_history])
        response = get_chatbot_response(chat_history_string, message)

        chat_history.append({'USER': message, 'AI': response})

        # Store the updated chat history in the session
        request.session['chat_history'] = chat_history

        return JsonResponse({'message': message, 'response': response})
    return render(request, "chatbot.html")


def get_chatbot_response(chat_history: str, request: str):
    """
    A function to generate a chatbot response to the user message
    :param chat_history: The previous chat history (conversation history)
    :param request: The request containing the user message
    :return: Generated chatbot response
    """
    response = ""
    try:
        response = client.chat.completions.create(
            model=config.ENGINE,
            messages=[
                {'role': 'user',
                 'content': f'{config.WEBSITE_DESCRIPTION_FOR_CHATBOT}.\n This is your chat history: {chat_history}. \nUSER:{request}\nYour Response:'}
            ],

        )
        final_response = response.choices[0].message.content

        logging.info(final_response)
        return final_response

    except Exception as e:
        logging.info(response.choices[0].message.content)
        logging.info("something went wrong with generating exercises based on the provided prompt.")
        logging.error(e)


@authenticated_user
def generate_adapted_content_view(request):
    """
    A view to generate an adapted content
    :param request: Request containing user message
    :return: Adapted content
    """
    if request.method == "POST":
        input_option = request.POST.get("input_option")
        if input_option == "write":
            input_text = request.POST.get("input_text")
            input_user_group = request.POST.get("input_user_group")
            if input_text:
                adapted_content = generate_adapted_content(input_text, input_user_group)

                return render(request, 'adapted_content_generation.html',
                                          {'adapted_content': adapted_content,
                                           "input_option": input_option})
            else:
                messages.error(request, "Please enter text that needs to be adapted for the target user group")
                return render(request, "adapted_content_generation.html")
        elif input_option == "upload":
            uploaded_file = request.FILES.get("file")
            input_user_group = request.POST.get("input_user_group")
            if uploaded_file:
                file_extension = os.path.splitext(uploaded_file.name)[1].lower()

                if file_extension in [".docx", ".pdf", ".pptx"]:
                    if file_extension == ".docx":
                        doc = docx.Document(uploaded_file)
                        full_text = []
                        for paragraph in doc.paragraphs:
                            full_text.append(paragraph.text)

                        document_text = "\n".join(full_text)
                        if document_text:
                            adapted_content = generate_adapted_content(document_text, input_user_group)

                            document = docx.Document()

                            document.add_heading(f"Document for: {input_user_group}", level=1)
                            document.add_paragraph(str(adapted_content))

                            response = HttpResponse(
                                content_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document')
                            response['Content-Disposition'] = 'attachment; filename="generated_document.docx"'

                            document.save(response)
                            return response
                        else:
                            messages.error(request, "You uploaded an empty .docx file.")

                    elif file_extension == ".pdf":
                        pdf_content = uploaded_file
                        pdf_file = PyPDF2.PdfReader(pdf_content)
                        pdf_text = ""

                        for page_num in range(len(pdf_file.pages)):
                            page = pdf_file.pages[page_num]
                            pdf_text += page.extract_text()

                        if pdf_text:
                            adapted_content = generate_adapted_content(pdf_text, input_user_group)

                            response = HttpResponse(content_type='application/pdf')
                            response['Content-Disposition'] = 'attachment; filename="generated_document.pdf"'

                            pdf_buffer = response

                            styles = getSampleStyleSheet()

                            story = []

                            title = f"Document for: {input_user_group}"
                            story.append(Paragraph(title, styles['Title']))

                            story.append(Spacer(1, 12))

                            content_paragraphs = adapted_content.split('\n')
                            for paragraph in content_paragraphs:
                                story.append(Paragraph(paragraph, styles['Normal']))

                            doc = SimpleDocTemplate(pdf_buffer, pagesize=letter)
                            doc.build(story)
                            return response
                        else:
                            messages.error(request,
                                            "You uploaded an empty .pdf file.")

                    elif file_extension == ".pptx":
                        presentation = Presentation(uploaded_file)

                        slide_text = []

                        for slide in presentation.slides:
                            slide_text.append("")
                            for shape in slide.shapes:
                                if hasattr(shape, "text"):
                                    slide_text.append(shape.text)

                        presentation_text = "\n".join(slide_text)
                        if presentation_text:
                            uploaded_content_summary = generate_summary(presentation_text)
                            adapted_content = generate_adapted_content(uploaded_content_summary, input_user_group)

                            presentation = generate_presentation(
                                f"Generate a presentation using the following content. The content is designed for: {input_user_group}. Do not forget to use images when appropriate to make the presentation more entertaining. Your images must be illustrative for the user group. The content is the following:\n" + str(adapted_content))
                            response = HttpResponse(
                                content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
                            response['Content-Disposition'] = 'attachment; filename="generated_presentation.pptx"'
                            presentation.save(response)

                            return response
                        else:
                            messages.error(request,
                                            "You uploaded an empty .pptx file.")
                else:
                    messages.error(request,
                                    "The Uploaded file must have one of the following extensions: .docs, .pptx, .txt, .pdf")

            else:
                messages.error(request, "Please upload a file")

    return render(request, "adapted_content_generation.html")


def loading_page_view(request):
    """
    A view to navigate to the loading page
    :param request: Request object
    :return: None (Navigates to the loading page)
    """
    return render(request, "loading_page.html", {})


def presentation_download(request):
    """
    A function that triggers download of the generated presentation
    :param request: Request object
    :return: None (the necessary file is downloaded)
    """

    presentation_id = cache.get("generated_presentation_id", "000")
    filename = f'generated_presentation_{presentation_id}.pptx'
    if filename:
        fs = FileSystemStorage(location=presentation_root + "/presentations/")
        if fs.exists(filename):
            response = FileResponse(fs.open(filename, 'rb'), content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
            response['Content-Disposition'] = f'attachment; filename="{filename}"'
            #fs.delete(filename)
            # fs.delete(filename)
            # cache.clear()

            return response

    messages.error(request, "No presentation was generated or the session has expired.")
    return render(request, "presentation_generation.html", {"input_option": "write"})


def presentation_status(request):
    """
    A function to check if the presentation is ready or not
    :param request: Request object
    :return: A json indicating if the presentation has been generated or not
    """

    presentation_id = cache.get("generated_presentation_id", "000")
    filename = f'generated_presentation_{presentation_id}.pdf'
    fs = FileSystemStorage(location=presentation_root + "/presentations/")

    if fs.exists(filename):

        return JsonResponse({'status': 'ready'})
    else:

        return JsonResponse({'status': 'pending'})


def exercise_status(request):
    """
    A function to check if the exercises are ready or not
    :param request: Request object
    :return: A json indicating if the exercises have been generated or not
    """
    exercise_id = cache.get("generated_exercise_id", "000")
    filename = f'generated_exercise_{exercise_id}.docx'
    fs = FileSystemStorage()

    if fs.exists(filename):
        return JsonResponse({'status': 'ready'})
    else:
        return JsonResponse({'status': 'pending'})


def presentation_preview(request):
    """
    A function that navigates to the presentation preview page
    :param request: Request object
    :return: None (redirects to the presentation preview page and displays the generated presentation)
    """
    presentation_id = cache.get("generated_presentation_id", "000")
    filename = f'generated_presentation_{presentation_id}.pdf'  # Assuming conversion to PDF is done.

    fs = FileSystemStorage(location=presentation_root + "/presentations/")
    presentation_obj_filename = f'generated_presentation_{presentation_id}.pptx'
    if fs.exists(filename):
        presentation_obj = get_object_or_404(Presentations, presentation = presentation_obj_filename)
        # pdf_url = fs.url(filename)
        pdf_url = '/' + fs.base_url.lstrip('/') + filename

        return render(request, "presentation_preview.html", {"presentation_pdf_url": pdf_url, "presentation_obj": presentation_obj})
    else:
        messages.error(request, "No presentation was ready for preview.")
        return redirect('generate_presentation')


def view_pdf(request):
    """
    A function that when called, returns the contents of the generated presentation pdf
    :param request:
    :return:
    """
    # Path to the PDF file
    # pdf_path = 'my_presentation.pdf'
    presentation_id = cache.get("generated_presentation_id", "000")
    pdf_path = f'generated_presentation_{presentation_id}.pdf'

    fs = FileSystemStorage(location=presentation_root + "/presentations/")
    if fs.exists(pdf_path):
        full_pdf_path = presentation_root + "/presentations/" + pdf_path
        with open(full_pdf_path, 'rb') as f:
            pdf_data = f.read()

    # Return the PDF content as the HTTP response
        fs.delete(pdf_path)
        return HttpResponse(pdf_data, content_type='application/pdf')
    else:
        messages.error(request, "Something went wrong while generating your presentation, please try again.")
        return redirect('generate_presentation')


def handle_modification_message(request):
    """
    A function to handle user message requiring a modification in the presentation
    :param request: Request containing user message
    :return: Appropriate response to the user message
    """
    logging.info("IN HANDLE MODIFICATION")
    if request.method == 'POST':
        data = json.loads(request.body)
        user_message = data.get('message', '')

        # Call your function to check user message
        modify_presentation = check_user_message(user_message)

        if modify_presentation:
            # Call your modify_presentation function here
            # Example: modify_presentation(user_message)
            response_data = {'modifyPresentation': True}
        else:
            # Call your function to answer user message
            # Example: response = answer_user_message(user_message)
            modification_assistant_answer = generate_modification_assistant_response(user_message)
            response_data = {'modifyPresentation': False, 'response': modification_assistant_answer}

        return JsonResponse(response_data)


def modify_presentation(request):
    """
    A function to modify the given presentation
    :param request: Request object
    :return: The modified presentation contents
    """
    logging.info("IN MODIFY PRESENTATION")
    if request.method == 'POST':
        data = json.loads(request.body)
        user_message = data.get('message', '')
        presentation_id = cache.get("generated_presentation_id")
        if not presentation_id:
            messages.error(request,
                           "Your session has expired")
            return render(request, "presentation_generation.html", {"input_option": "write"})
        generated_presentation_json = cache.get("generated_presentation_json", "{}")
        logging.info(f"GENERATED PRESENTATION JSON: {generated_presentation_json}")
        generated_presentation_json = json.loads(generated_presentation_json)
        modified_presentation_json, modified_presentation_object = generate_modified_presentation(generated_presentation_json, user_message)
        modified_presentation_info = get_presentation_info(modified_presentation_json)
        cache.set("generated_presentation_json", json.dumps(modified_presentation_json), timeout=1500)
        presentation_id = cache.get("generated_presentation_id", "000")
        filename = f'generated_presentation_{presentation_id}.pptx'
        fs = FileSystemStorage(location=presentation_root + "/presentations/")
        if fs.exists(filename):
            pass
        with fs.open(filename, 'wb') as pptx_file:
            modified_presentation_object.save(pptx_file)

        if fs.exists(filename):
            pptx_file_path = fs.path(filename)
            convert_pptx_to_pdf(pptx_file_path)

        #Update the existing presentation object with the new modified information
        if modified_presentation_info:
            presentation_obj = get_object_or_404(Presentations, presentation=filename)
            presentation_obj.main_title = modified_presentation_info['main_title']
            presentation_obj.titles = modified_presentation_info['titles']
            presentation_obj.save(update_fields=['main_title', 'titles'])

        pdf_path = f'generated_presentation_{presentation_id}.pdf'
        pdf_full_path = presentation_root + "/presentations/" + pdf_path
        if fs.exists(pdf_path):
            with open(pdf_full_path, 'rb') as f:
                pdf_data = f.read()

            # Return the PDF content as the HTTP response
            fs.delete(pdf_path)
            return HttpResponse(pdf_data, content_type='application/pdf')
        else:
            messages.error(request, "Something went wrong while generating your modified presentation, please try again.")
            return render(request, "presentation_generation.html", {"input_option": "write"})


def exercise_loading_page_view(request):
    """
    A function to redirect to the exercise loading page
    A function to redirect to the exercise loading page
    :param request: Request object
    :return: None (redirects to the exercise loading page)
    """
    return render(request, "exercise_loading_page.html", {})


def get_exercise_view(request):
    """
    A view that generates exercises and redirects to the exercise generation page
    :param request: Request containing user message
    :return: None (redirects to the exercise page and displays the generated exercises)
    """
    exercise_id = cache.get("generated_exercise_id", "000")
    filename = f'generated_exercise_{exercise_id}.docx'
    # input_option = request.session.get('input_option')
    input_option = cache.get("input_option", "invalid")

    generated_exercises = ""
    fs = FileSystemStorage()
    if fs.exists(filename):
        with fs.open(filename, 'rb') as docx_file:
            doc = Document(docx_file)
            generated_exercises = []

            for paragraph in doc.paragraphs:
                generated_exercises.append(paragraph.text)
    fs.delete(filename)

    return render(request, 'exercise_generation.html',
                    {'generated_exercises': generated_exercises[0],
                    "input_option": input_option})
