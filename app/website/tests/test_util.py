import json
from app.util.generate_presentation import get_presentation_json, process_and_store_presentation_json
from django.test import TestCase
from pptx import *

class TestPresentationCreation(TestCase):
    # def test_create_presentation(self):
    #     # Define input parameters
    #     input_param1 = "A presentation on Linear algebra"

    #     # Call the function being tested
    #     result = get_presentation_json(input_param1)

    #     json_object = json.dumps(result, indent = 4)  

    #     # Ensure the result is a valid JSON string
    #     assert json.loads(json_object)

    def test_process_and_store_presentation_json(self):

        # Sample input JSON
        input_json = {'slide1': {'slide_layout': 0, 'title': 'Title Slide', 'subtitle': 'Subtitle Text', 'space_after_paragraphs': -1, 'paragraphs': {}}, 
             'slide2': {'slide_layout': 1, 'title': 'Slide1', 'subtitle': '', 'space_after_paragraphs': 0.1, 'paragraphs': {'paragraph1': 'Paragraph 1 Text'}}, 
             'slide3': {'slide_layout': 1, 'title': 'Slide2', 'subtitle': '', 'space_after_paragraphs': 0.1, 'paragraphs': {'paragraph1': 'Paragraph 1 Text'}},
        }
        mock_presentation = Presentation()

        # Call the function with the sample input
        presentation = process_and_store_presentation_json(input_json)

        # Get the object type for the functions
        function_return_type = type(presentation)
        comparison_type = type(mock_presentation)

        #compare data object types
        self.assertTrue(function_return_type == comparison_type)

        #for loop to check the values in every slide are correct

        error = 0
        for slide in presentation.slides:
            
            if slide.shapes.title.text == "Title Slide":
                if slide.shapes.placeholders[1].text == "Subtitle Text":
                    pass
                else:
                    error = 1

            elif slide.shapes.title.text == "Slide1":
                if slide.shapes.placeholders[1].text[1:] == "Paragraph 1 Text":
                    pass
                else:
                    error = 1

            elif slide.shapes.title.text == "Slide2":
                if slide.shapes.placeholders[1].text[1:] == "Paragraph 1 Text":
                    pass
                else:
                    error = 1
            else:
                error = 1

        self.assertTrue(error == 0)







    

