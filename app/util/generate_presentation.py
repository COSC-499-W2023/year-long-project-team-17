from pptx import Presentation
from pptx.util import Pt, Inches
import tiktoken
import os
from typing import Dict
import openai
from openai import OpenAI
import json
import requests
from dotenv import load_dotenv, find_dotenv
import logging
from . import config
from io import BytesIO
from PIL import Image
import random
from pptx.shapes.shapetree import PicturePlaceholder, SlidePlaceholder
SlidePlaceholder.insert_picture = PicturePlaceholder.insert_picture
SlidePlaceholder._new_placeholder_pic = PicturePlaceholder._new_placeholder_pic
SlidePlaceholder._get_or_add_image = PicturePlaceholder._get_or_add_image

logging.basicConfig(level=logging.DEBUG)

load_dotenv(find_dotenv())

client = OpenAI(
    api_key=os.environ.get("OPENAI_API_KEY"),
)


def get_presentation_json(description: str) -> Dict[str, str]:
    """
        Generates Presentation json based on content.
        Args:
            description (str): A description of the user's presentation.
        Returns:
             presentation in json format
        """
    response = ""
    try:

        response = client.chat.completions.create(
            model=config.ENGINE,
            messages=[
                {'role': 'user',
                 'content': f'{config.PRESENTATION_GENERATION_PROMPT}\nUSER:{description}\nYour Response:'}
            ],
            response_format={"type": "json_object"},
        )
        response_json = json.loads(response.choices[0].message.content)

        logging.info(response_json)
        return response_json
    except Exception as e:
        logging.info(response.choices[0].message.content)

        logging.info("something went wrong with generating presentation json based on the provided content.")
        logging.info(e)
        logging.error(e)


def download_image(url, destination):
    response = requests.get(url)
    if response.status_code == 200:
        original_image = Image.open(BytesIO(response.content))

        original_image.save(destination)

        logging.info(f"Image downloaded successfully and saved to {destination}")
    else:
        logging.info(f"Failed to download image. Status code: {response.status_code}")


def process_and_store_presentation_json(result: dict, modified=False):
    """
           Process presentation json and create a Presentation object
           Args:
               result (dict): The json of presentation
           Returns:
                The presentation as a Presentation object
           """

    try:
        presentation_template_index = random.randint(1, 10)
        # path = "../media/presentation_templates/presentation_template_1.pptx"
        presentation_template_path = f"app/media/presentation_templates/presentation_template_{presentation_template_index}.pptx"
        print(presentation_template_path)
        presentation = Presentation(presentation_template_path)
        i = 0
        all_image_paths = []
        for key, value in result.items():
            # print(key, ":", value)
            current_slide = presentation.slides.add_slide(presentation.slide_layouts[value["slide_layout"]])
            current_slide_title = current_slide.shapes.title
            current_slide_title.text = value["title"]
            if value.get("subtitle"):
                current_slide_subtitle = current_slide.shapes.placeholders[1]
                current_slide_subtitle.text = value["subtitle"]
            if value.get("paragraphs"):
                for paragraph, paragraph_text in value["paragraphs"].items():
                    current_content = None
                    if value["slide_layout"] == 3:
                        # print("true")
                        current_content = current_slide.shapes.placeholders[1]
                    else:
                        current_content = current_slide.shapes.placeholders[1]
                    current_paragraph = current_content.text_frame.add_paragraph()
                    current_paragraph.text = paragraph_text

                    for run in current_paragraph.runs:
                        run.font.size = Pt(18)

                    if value.get("space_after_paragraphs") > 0:
                        current_paragraph.space_after = Inches(value["space_after_paragraphs"])
            if value.get("image_title") and value.get("slide_layout") == 3:
                # content = current_slide.shapes.placeholders[2]
                # image_title = content.text_frame.add_paragraph()
                # image_title.text = value.get("image_title")
                # for run in image_title.runs:
                #     run.font.size = Pt(18)
                image_path = generate_presentation_images(value.get("image_keywords"), i)
                all_image_paths.append(image_path)
                content_placeholder = current_slide.shapes.placeholders[2]
                content_placeholder.insert_picture(image_path)
                # left_inch = Inches(5.5)
                # top_inch = Inches(3)
                # width_inch = height_inch = Inches(4)
                # current_slide.shapes.add_picture(image_path, left_inch, top_inch, width_inch, height_inch)
                i += 1
        if modified:
            presentation.save("my_presentation_modified.pptx")
        else:
            presentation.save("my_presentation.pptx")

        for image_path in all_image_paths:
            if os.path.exists(image_path):
                os.remove(image_path)
        return presentation
    except Exception as e:

        logging.info("something went wrong with generating Presentation object based on the presentation json.")
        logging.info(e)
        logging.error(e)

def get_presentation_info(result: dict):
    pres_info = {}
    pres_info['main_title'] = ""
    pres_info['titles'] = ""
    for i,value in enumerate(result.values()):
        #Stores the title in the first slide
        if i == 0:
            pres_info['main_title'] = value['title']
        #stores the titles in the remaining slides 
        elif value['title'].lower() != "conclusion" and value['title'].lower() != "references" and value['title'].lower() != "thank you":
            pres_info['titles'] += value['title'] + ", "
    titles = pres_info['titles']
    #remove space and comma at the end of titles string
    pres_info['titles'] = titles[:-2]
    return pres_info


def generate_presentation(description: str):
    """
    A wrapper function that generates a stores presentation end-to-end
    :param description: Description of the topic on what presentation needs to be generated
    :return: Generated presentation
    """
    try:
        presentation_json = get_presentation_json(description)
        pres_info = get_presentation_info(presentation_json)
        presentation_object = process_and_store_presentation_json(presentation_json, False)

        values = {}
        values['presentation'] = presentation_object
        values['pres_info'] = pres_info

        return values, presentation_json
    except Exception as e:
        logging.info("Something went wrong with generating presentation.")
        logging.info(e)
        logging.error(e)


def generate_presentation_images(image_keywords, image_number):
    """
    A function to generate images for the presentation using Dall-E
    :param image_keywords: Keywords that need to be used for image generation
    :param image_number: The index of the image
    :return: Generated image
    """
    result = client.images.generate(
        model="dall-e-3",
        prompt=image_keywords,
        size="1024x1024",
        quality="standard",
        n=1,
    )

    image_url = result.data[0].url
    image_path = f"downloaded_image_{image_number}.png"
    download_image(image_url, image_path)
    return image_path
