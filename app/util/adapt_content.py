from openai import OpenAI
from dotenv import find_dotenv, load_dotenv
import os
from . import config
import logging
from pptx import Presentation
import json
from pptx.util import Pt, Inches
# from generate_presentation import generate_presentation
logging.basicConfig(level=logging.DEBUG)

load_dotenv(find_dotenv())

client = OpenAI(
    api_key=os.environ.get("OPENAI_API_KEY"),
)


def extract_presentation_content(pptx_path):
    presentation = Presentation(pptx_path)

    presentation_content = ""
    for i, slide in enumerate(presentation.slides):
        presentation_content = presentation_content + f"Slide {i + 1}:\n"

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                presentation_content = presentation_content + (f"-{shape.text}\n")

    logging.info(f"Presentation Content: {presentation_content}")
    return presentation_content


def generate_adapted_content(original_content: str, target_user_group: str) -> str:
    """

    :param original_content: Original content provided by the user
    :param target_user_group: The target user group for which the content has to be adapted
    :return: The adapted content
    """
    adapted_content = ""
    try:
        logging.info("Starting to generate adapted content")
        adapted_content = client.chat.completions.create(
                    model=config.CONTENT_ADAPTATION_ENGINE,
                    messages=[
                        {'role': 'user',
                         'content': f'{config.CONTENT_ADAPTATION_SYSTEM_PROMPT}\n```USER\nTarget User Group: {target_user_group}\nContent:{original_content}\nYour Response:'}
                    ],
                )

        logging.info("Adapted presentation content")
        logging.info(adapted_content)
        return adapted_content.choices[0].message.content
    except Exception as e:
        logging.error(f"Something went wrong with generating adapted materials: {e}")
        return adapted_content


if __name__ == "__main__":
    original_content = """Multivariable calculus is a branch of mathematics that extends the principles of calculus to functions of several variables. Unlike single-variable calculus, which primarily deals with functions of a single independent variable, multivariable calculus involves functions with multiple independent variables. This field is essential for understanding and analyzing complex phenomena in various scientific and engineering disciplines, such as physics, economics, and computer science.

One key concept in multivariable calculus is the partial derivative, which measures the rate of change of a function with respect to one of its variables while keeping the other variables constant. Partial derivatives enable the analysis of how a function responds to changes in specific directions, providing a more comprehensive understanding of its behavior. Another fundamental tool is the gradient vector, which encapsulates the partial derivatives of a function and points in the direction of the steepest ascent. This vector plays a crucial role in optimization problems, helping to find the maximum or minimum values of a multivariable function.

Integral calculus also extends to multiple dimensions in multivariable calculus. Instead of integrating over a one-dimensional interval, the double and triple integrals involve integrating over regions in two or three-dimensional space, respectively. These integrals have applications in computing areas, volumes, and various physical quantities. Stokes' theorem and the divergence theorem are advanced theorems that relate integrals over surfaces and volumes, providing powerful tools for solving problems in physics, fluid dynamics, and electromagnetism. In summary, multivariable calculus is a rich and versatile field that enables a deeper understanding of complex systems and is indispensable in many scientific and engineering applications."""
    target_user_group = "first grade elementary school students"
    adapted_content = generate_adapted_content(original_content, target_user_group)
